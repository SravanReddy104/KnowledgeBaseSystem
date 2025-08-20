import os
from typing import List
from bs4 import BeautifulSoup
import pandas as pd
from markdown import markdown

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from langchain_community.document_loaders import PyPDFLoader, TextLoader

try:
    from docx import Document as DocxDocument
except Exception:  # pragma: no cover
    DocxDocument = None

try:
    from pptx import Presentation  # python-pptx
except Exception:  # pragma: no cover
    Presentation = None


def ensure_dir(path: str) -> None:
    os.makedirs(path, exist_ok=True)


def _load_pptx(path: str) -> List[Document]:
    if Presentation is None:
        raise RuntimeError("python-pptx not available to read .pptx files")
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    text = "\n".join(t for t in texts if t)
    return [Document(page_content=text, metadata={"source": os.path.basename(path)})]


def _load_html(path: str) -> List[Document]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    text = soup.get_text("\n")
    return [Document(page_content=text, metadata={"source": os.path.basename(path)})]


def load_single_file(path: str) -> List[Document]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return PyPDFLoader(path).load()
    if ext in {".txt", ".log"}:
        return TextLoader(path, encoding="utf-8").load()
    if ext == ".pptx":
        return _load_pptx(path)
    if ext in {".html", ".htm"}:
        return _load_html(path)
    if ext == ".md":
        # Convert Markdown to HTML then plain text
        with open(path, "r", encoding="utf-8") as f:
            html = markdown(f.read())
        text = BeautifulSoup(html, "html.parser").get_text("\n")
        return [Document(page_content=text, metadata={"source": os.path.basename(path)})]
    if ext == ".csv":
        df = pd.read_csv(path)
        text = df.to_markdown(index=False)
        return [Document(page_content=text, metadata={"source": os.path.basename(path)})]
    if ext == ".docx":
        if DocxDocument is None:
            raise RuntimeError("python-docx not available to read .docx files")
        d = DocxDocument(path)
        text = "\n".join(p.text for p in d.paragraphs)
        return [Document(page_content=text, metadata={"source": os.path.basename(path)})]

    raise ValueError(f"Unsupported file type: {ext}")


def load_files(paths: List[str]) -> List[Document]:
    docs: List[Document] = []
    for p in paths:
        docs.extend(load_single_file(p))
    return docs


def chunk_documents(docs: List[Document], chunk_size: int = 1200, chunk_overlap: int = 200) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", " ", ""],
    )
    return splitter.split_documents(docs)
